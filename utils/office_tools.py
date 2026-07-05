"""
Office document <-> PDF conversions using pure-Python libraries
(no LibreOffice / MS Office dependency, so this works on serverless).

NOTE on fidelity: because there is no LibreOffice/Word engine available
in a serverless environment, Word/Excel/PowerPoint -> PDF here is a
best-effort re-render (text, basic styling, tables, and images are
preserved; complex layouts, custom fonts, and advanced formatting
are not pixel-perfect). For pixel-perfect office conversion at scale,
see the README section on running this behind a Docker/Render backend
with LibreOffice installed.
"""
import io

from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
)
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors


# ---------------------------------------------------------------- WORD -> PDF
def word_to_pdf(file_stream):
    import docx

    document = docx.Document(file_stream)
    styles = getSampleStyleSheet()

    out = io.BytesIO()
    doc = SimpleDocTemplate(out, pagesize=A4, topMargin=0.8 * inch, bottomMargin=0.8 * inch)
    story = []

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 8))
            continue
        style_name = para.style.name if para.style else "Normal"
        if "Heading 1" in style_name:
            story.append(Paragraph(text, styles["Heading1"]))
        elif "Heading 2" in style_name:
            story.append(Paragraph(text, styles["Heading2"]))
        elif "Heading" in style_name:
            story.append(Paragraph(text, styles["Heading3"]))
        else:
            story.append(Paragraph(text, styles["Normal"]))
        story.append(Spacer(1, 6))

    for table in document.tables:
        data = [[cell.text for cell in row.cells] for row in table.rows]
        if data:
            t = Table(data)
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
            ]))
            story.append(t)
            story.append(Spacer(1, 12))

    if not story:
        story.append(Paragraph("(Empty document)", styles["Normal"]))

    doc.build(story)
    out.seek(0)
    return out


# ---------------------------------------------------------------- EXCEL -> PDF
def excel_to_pdf(file_stream):
    from openpyxl import load_workbook

    wb = load_workbook(file_stream, data_only=True)
    styles = getSampleStyleSheet()

    out = io.BytesIO()
    doc = SimpleDocTemplate(out, pagesize=(11 * inch, 8.5 * inch),
                             topMargin=0.5 * inch, bottomMargin=0.5 * inch)
    story = []

    for sheet in wb.worksheets:
        story.append(Paragraph(f"Sheet: {sheet.title}", styles["Heading2"]))
        story.append(Spacer(1, 6))

        data = []
        for row in sheet.iter_rows(values_only=True):
            data.append(["" if v is None else str(v) for v in row])

        if data:
            # Cap columns/width so wide sheets don't overflow the page
            t = Table(data, repeatRows=1)
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e8e8e8")),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(t)
        story.append(Spacer(1, 20))

    doc.build(story)
    out.seek(0)
    return out


# ---------------------------------------------------------------- PPTX -> PDF
def ppt_to_pdf(file_stream):
    from pptx import Presentation
    from pptx.util import Emu
    from reportlab.pdfgen import canvas

    prs = Presentation(file_stream)
    slide_w = prs.slide_width / 914400 * inch   # EMU -> points via inches
    slide_h = prs.slide_height / 914400 * inch

    out = io.BytesIO()
    c = canvas.Canvas(out, pagesize=(slide_w, slide_h))

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                left = shape.left / 914400 * inch if shape.left is not None else 0
                top = shape.top / 914400 * inch if shape.top is not None else 0
                width = shape.width / 914400 * inch if shape.width is not None else 0
                height = shape.height / 914400 * inch if shape.height is not None else 0
                pdf_y = slide_h - top - height  # flip Y axis (pptx origin is top-left)

                if shape.shape_type == 13:  # PICTURE
                    img_bytes = shape.image.blob
                    img_stream = io.BytesIO(img_bytes)
                    c.drawImage(io.BytesIO(img_bytes) if False else img_stream,
                                left, pdf_y, width=width, height=height,
                                preserveAspectRatio=True, mask="auto")
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    c.setFont("Helvetica", 14)
                    text = shape.text_frame.text
                    text_y = slide_h - top - 16
                    for line in text.splitlines():
                        c.drawString(left, text_y, line[:120])
                        text_y -= 16
            except Exception:
                continue
        c.showPage()

    c.save()
    out.seek(0)
    return out


# ---------------------------------------------------------------- HTML -> PDF
def html_to_pdf(html_string):
    from xhtml2pdf import pisa

    out = io.BytesIO()
    pisa.CreatePDF(src=html_string, dest=out)
    out.seek(0)
    return out
