"""
Core PDF manipulation utilities.
All functions take/return in-memory bytes (BytesIO) so they work cleanly
in stateless serverless environments (Vercel /tmp is ephemeral).
"""
import io
import zipfile
import difflib

import fitz  # PyMuPDF
import pikepdf
from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4


# ---------------------------------------------------------------- MERGE
def merge_pdfs(file_streams):
    writer = PdfWriter()
    for stream in file_streams:
        reader = PdfReader(stream)
        for page in reader.pages:
            writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- SPLIT
def split_pdf(file_stream, ranges_str=None):
    """
    ranges_str: e.g. "1-3,5,7-9". If None/empty -> split into one PDF per page.
    Returns a ZIP file (BytesIO) containing the resulting PDFs.
    """
    reader = PdfReader(file_stream)
    n_pages = len(reader.pages)

    def parse_ranges(s):
        parts = []
        for chunk in s.split(","):
            chunk = chunk.strip()
            if not chunk:
                continue
            if "-" in chunk:
                a, b = chunk.split("-")
                parts.append((int(a) - 1, int(b) - 1))
            else:
                parts.append((int(chunk) - 1, int(chunk) - 1))
        return parts

    ranges = parse_ranges(ranges_str) if ranges_str else [(i, i) for i in range(n_pages)]

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for idx, (start, end) in enumerate(ranges, start=1):
            writer = PdfWriter()
            for p in range(start, min(end, n_pages - 1) + 1):
                writer.add_page(reader.pages[p])
            part_buf = io.BytesIO()
            writer.write(part_buf)
            zf.writestr(f"part_{idx}.pdf", part_buf.getvalue())
    zip_buf.seek(0)
    return zip_buf


# ---------------------------------------------------------------- COMPRESS
def compress_pdf(file_stream, level="recommended"):
    """
    level: 'low' (best quality), 'recommended', 'extreme' (smallest size)
    Uses PyMuPDF to downsample embedded images, then pikepdf to
    strip/compress streams and remove unused objects.
    """
    quality_map = {"low": 85, "recommended": 60, "extreme": 30}
    dpi_scale_map = {"low": 1.0, "recommended": 0.75, "extreme": 0.5}
    jpg_quality = quality_map.get(level, 60)
    scale = dpi_scale_map.get(level, 0.75)

    src = fitz.open(stream=file_stream.read(), filetype="pdf")
    for page in src:
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                base = src.extract_image(xref)
                pix = fitz.Pixmap(base["image"])
                if pix.n - pix.alpha >= 4:  # CMYK -> RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                if scale != 1.0:
                    new_w = max(1, int(pix.width * scale))
                    new_h = max(1, int(pix.height * scale))
                    pix = fitz.Pixmap(pix, new_w, new_h, None)
                new_bytes = pix.tobytes("jpeg", jpg_quality=jpg_quality)
                src.update_stream(xref, new_bytes, new=False)
            except Exception:
                continue  # skip images that can't be safely re-encoded

    mid = io.BytesIO()
    src.save(mid, garbage=4, deflate=True, clean=True)
    src.close()
    mid.seek(0)

    # Final pass with pikepdf for extra structural compression
    out = io.BytesIO()
    with pikepdf.open(mid) as pdf:
        pdf.save(out, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)
    out.seek(0)
    return out


# ---------------------------------------------------------------- ROTATE
def rotate_pdf(file_stream, angle=90, pages_str=None):
    reader = PdfReader(file_stream)
    writer = PdfWriter()
    n = len(reader.pages)
    target = set(range(n))
    if pages_str:
        target = set(int(p.strip()) - 1 for p in pages_str.split(",") if p.strip())
    for i, page in enumerate(reader.pages):
        if i in target:
            page.rotate(angle)
        writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- WATERMARK
def watermark_pdf(file_stream, text="CONFIDENTIAL", opacity=0.3, font_size=48):
    reader = PdfReader(file_stream)
    writer = PdfWriter()

    for page in reader.pages:
        w = float(page.mediabox.width)
        h = float(page.mediabox.height)

        wm_buf = io.BytesIO()
        c = canvas.Canvas(wm_buf, pagesize=(w, h))
        c.saveState()
        c.setFillAlpha(opacity)
        c.setFont("Helvetica-Bold", font_size)
        c.translate(w / 2, h / 2)
        c.rotate(45)
        c.drawCentredString(0, 0, text)
        c.restoreState()
        c.save()
        wm_buf.seek(0)

        wm_reader = PdfReader(wm_buf)
        page.merge_page(wm_reader.pages[0])
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- PAGE NUMBERS
def add_page_numbers(file_stream, position="bottom-center", start_at=1):
    reader = PdfReader(file_stream)
    writer = PdfWriter()
    n = len(reader.pages)

    positions = {
        "bottom-center": lambda w, h: (w / 2, 20),
        "bottom-right": lambda w, h: (w - 50, 20),
        "bottom-left": lambda w, h: (50, 20),
        "top-center": lambda w, h: (w / 2, h - 30),
    }
    pos_fn = positions.get(position, positions["bottom-center"])

    for i, page in enumerate(reader.pages):
        w = float(page.mediabox.width)
        h = float(page.mediabox.height)
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=(w, h))
        c.setFont("Helvetica", 10)
        x, y = pos_fn(w, h)
        c.drawCentredString(x, y, str(start_at + i))
        c.save()
        buf.seek(0)
        overlay = PdfReader(buf).pages[0]
        page.merge_page(overlay)
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- PROTECT / UNLOCK
def protect_pdf(file_stream, password):
    reader = PdfReader(file_stream)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


def unlock_pdf(file_stream, password):
    reader = PdfReader(file_stream)
    if reader.is_encrypted:
        reader.decrypt(password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- ORGANIZE (reorder/delete)
def organize_pdf(file_stream, order_str):
    """order_str e.g. '3,1,2' (1-indexed) - keeps only listed pages, in that order."""
    reader = PdfReader(file_stream)
    writer = PdfWriter()
    order = [int(x.strip()) - 1 for x in order_str.split(",") if x.strip()]
    for idx in order:
        if 0 <= idx < len(reader.pages):
            writer.add_page(reader.pages[idx])
    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- CROP
def crop_pdf(file_stream, left=0, bottom=0, right=0, top=0):
    """Margins in points to trim from each edge, applied to all pages."""
    reader = PdfReader(file_stream)
    writer = PdfWriter()
    for page in reader.pages:
        box = page.mediabox
        box.lower_left = (float(box.left) + left, float(box.bottom) + bottom)
        box.upper_right = (float(box.right) - right, float(box.top) - top)
        writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- REDACT (true content removal)
def redact_pdf(file_stream, terms):
    """terms: list of strings to find & permanently remove from the page content."""
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    for page in doc:
        for term in terms:
            areas = page.search_for(term)
            for rect in areas:
                page.add_redact_annot(rect, fill=(0, 0, 0))
        page.apply_redactions()
    out = io.BytesIO()
    doc.save(out, garbage=4, deflate=True)
    doc.close()
    out.seek(0)
    return out


# ---------------------------------------------------------------- REPAIR
def repair_pdf(file_stream):
    out = io.BytesIO()
    with pikepdf.open(file_stream) as pdf:
        pdf.save(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- COMPARE (text diff)
def compare_pdfs(file_stream_a, file_stream_b):
    import pdfplumber

    def extract_text(stream):
        text_pages = []
        with pdfplumber.open(stream) as pdf:
            for page in pdf.pages:
                text_pages.append(page.extract_text() or "")
        return "\n".join(text_pages)

    text_a = extract_text(file_stream_a).splitlines()
    text_b = extract_text(file_stream_b).splitlines()

    diff = difflib.HtmlDiff(wrapcolumn=80).make_table(
        text_a, text_b, fromdesc="Document A", todesc="Document B", context=True, numlines=2
    )
    return diff


# ---------------------------------------------------------------- PDF -> JPG
def pdf_to_jpg(file_stream, dpi=150):
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i}.jpg", pix.tobytes("jpeg"))
    doc.close()
    zip_buf.seek(0)
    return zip_buf


# ---------------------------------------------------------------- PDF -> WORD
def pdf_to_word(file_stream):
    import tempfile, os
    from pdf2docx import Converter

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_in:
        tmp_in.write(file_stream.read())
        tmp_in_path = tmp_in.name
    tmp_out_path = tmp_in_path.replace(".pdf", ".docx")

    try:
        cv = Converter(tmp_in_path)
        cv.convert(tmp_out_path)
        cv.close()
        with open(tmp_out_path, "rb") as f:
            out = io.BytesIO(f.read())
        out.seek(0)
        return out
    finally:
        for p in (tmp_in_path, tmp_out_path):
            if os.path.exists(p):
                os.remove(p)


# ---------------------------------------------------------------- PDF -> PPTX (image-based, preserves visuals)
def pdf_to_ppt(file_stream, dpi=150):
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)

    prs = Presentation()
    for page in doc:
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")

        # slide size matches page aspect ratio (EMU: 914400 per inch)
        page_w_in = page.rect.width / 72
        page_h_in = page.rect.height / 72
        prs.slide_width = Emu(int(page_w_in * 914400))
        prs.slide_height = Emu(int(page_h_in * 914400))

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

    doc.close()
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ---------------------------------------------------------------- PDF -> EXCEL (table extraction)
def pdf_to_excel(file_stream):
    import pdfplumber
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)

    with pdfplumber.open(file_stream) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables()
            if not tables:
                continue
            for t_idx, table in enumerate(tables, start=1):
                sheet_name = f"Page{i}_T{t_idx}"[:31]
                ws = wb.create_sheet(title=sheet_name)
                for row in table:
                    ws.append([cell if cell is not None else "" for cell in row])

        if len(wb.sheetnames) == 0:
            # fallback: no tables found -> dump raw text per page
            for i, page in enumerate(pdf.pages, start=1):
                ws = wb.create_sheet(title=f"Page{i}")
                text = page.extract_text() or ""
                for line_no, line in enumerate(text.splitlines(), start=1):
                    ws.cell(row=line_no, column=1, value=line)

    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    return out
